"""
Document parser service.
Supports: PDF, DOCX, DOC, XLSX, XLS, PPTX, PPT, TXT, HTML, Markdown
"""
import logging
import io
from dataclasses import dataclass, field
from typing import List, Optional

logger = logging.getLogger(__name__)


@dataclass
class ParsedPage:
    page: int
    text: str
    meta: dict = field(default_factory=dict)


@dataclass
class ParseResult:
    text: str
    pages: List[ParsedPage]
    page_count: int
    word_count: int
    language: str = "en"
    meta: dict = field(default_factory=dict)


class DocumentParser:

    def parse_bytes(self, data: bytes, mime_type: str, filename: str = "") -> ParseResult:
        ext = filename.rsplit(".", 1)[-1].lower() if "." in filename else ""

        if mime_type == "application/pdf" or ext == "pdf":
            return self._parse_pdf(data)
        elif mime_type in (
            "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
            "application/msword",
        ) or ext in ("docx", "doc"):
            return self._parse_docx(data)
        elif mime_type in (
            "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
            "application/vnd.ms-excel",
        ) or ext in ("xlsx", "xls"):
            return self._parse_xlsx(data)
        elif mime_type in (
            "application/vnd.openxmlformats-officedocument.presentationml.presentation",
            "application/vnd.ms-powerpoint",
        ) or ext in ("pptx", "ppt"):
            return self._parse_pptx(data)
        elif mime_type == "text/html" or ext == "html":
            return self._parse_html(data.decode("utf-8", errors="ignore"))
        elif mime_type == "text/markdown" or ext in ("md", "markdown"):
            return self._parse_markdown(data.decode("utf-8", errors="ignore"))
        elif mime_type == "text/plain" or ext == "txt":
            return self._parse_text(data.decode("utf-8", errors="ignore"))
        else:
            # Fallback: try as plain text
            return self._parse_text(data.decode("utf-8", errors="ignore"))

    def _parse_pdf(self, data: bytes) -> ParseResult:
        import pypdf
        reader = pypdf.PdfReader(io.BytesIO(data))
        pages = []
        full_text_parts = []
        for i, page in enumerate(reader.pages):
            text = page.extract_text() or ""
            pages.append(ParsedPage(page=i + 1, text=text))
            full_text_parts.append(text)
        full_text = "\n\n".join(full_text_parts)
        return ParseResult(
            text=full_text,
            pages=pages,
            page_count=len(pages),
            word_count=len(full_text.split()),
        )

    def _parse_docx(self, data: bytes) -> ParseResult:
        from docx import Document
        doc = Document(io.BytesIO(data))
        paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
        full_text = "\n\n".join(paragraphs)
        pages = [ParsedPage(page=1, text=full_text)]
        return ParseResult(
            text=full_text,
            pages=pages,
            page_count=1,
            word_count=len(full_text.split()),
        )

    def _parse_xlsx(self, data: bytes) -> ParseResult:
        import openpyxl
        wb = openpyxl.load_workbook(io.BytesIO(data), read_only=True, data_only=True)
        pages = []
        full_parts = []
        for i, sheet in enumerate(wb.worksheets):
            rows = []
            for row in sheet.iter_rows(values_only=True):
                row_text = "\t".join(str(c) if c is not None else "" for c in row)
                if row_text.strip():
                    rows.append(row_text)
            sheet_text = f"Sheet: {sheet.title}\n" + "\n".join(rows)
            pages.append(ParsedPage(page=i + 1, text=sheet_text))
            full_parts.append(sheet_text)
        full_text = "\n\n".join(full_parts)
        return ParseResult(
            text=full_text,
            pages=pages,
            page_count=len(pages),
            word_count=len(full_text.split()),
        )

    def _parse_pptx(self, data: bytes) -> ParseResult:
        from pptx import Presentation
        prs = Presentation(io.BytesIO(data))
        pages = []
        full_parts = []
        for i, slide in enumerate(prs.slides):
            texts = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    texts.append(shape.text)
            slide_text = "\n".join(texts)
            pages.append(ParsedPage(page=i + 1, text=slide_text))
            full_parts.append(slide_text)
        full_text = "\n\n".join(full_parts)
        return ParseResult(
            text=full_text,
            pages=pages,
            page_count=len(pages),
            word_count=len(full_text.split()),
        )

    def _parse_html(self, html: str) -> ParseResult:
        from bs4 import BeautifulSoup
        soup = BeautifulSoup(html, "lxml")
        # Remove scripts and styles
        for tag in soup(["script", "style", "nav", "footer", "header"]):
            tag.decompose()
        full_text = soup.get_text(separator="\n", strip=True)
        pages = [ParsedPage(page=1, text=full_text)]
        return ParseResult(
            text=full_text,
            pages=pages,
            page_count=1,
            word_count=len(full_text.split()),
        )

    def _parse_markdown(self, text: str) -> ParseResult:
        import markdown2
        html = markdown2.markdown(text)
        return self._parse_html(html)

    def _parse_text(self, text: str) -> ParseResult:
        pages = [ParsedPage(page=1, text=text)]
        return ParseResult(
            text=text,
            pages=pages,
            page_count=1,
            word_count=len(text.split()),
        )

